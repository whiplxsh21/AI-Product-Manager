import base64
import os
import re
import tempfile
from datetime import datetime

from langchain_core.messages import HumanMessage, SystemMessage

from runtime import effective_config
from database import SessionLocal
from llm import get_ingestion_llm, get_vision_llm, invoke_with_retry
from models import Document, GeneratedOutput, PipelineRun, Project
from pipeline.prompts import CLEANING_PROMPT, VISION_PROMPT
from pipeline.state import PipelineState
from storage import get_storage

# Max words per cleaning request. A cleaning call's tokens = input + output, and
# output ≈ input, so ~1500 words (~2k input + ~2k output tokens) stays under the
# tightest free-tier budget (Groq 8b-instant = 6000 tokens/min). Smaller
# transcripts are cleaned in a single call; larger ones split into a few calls,
# each cached per document so the cost is paid only once.
_MAX_SINGLE_PASS_WORDS = 1500

# Conservative filler tokens removed by the regex cleaner. Kept deliberately
# small to avoid stripping substantive words — the framework LLM tolerates
# residual noise, so the goal is just to thin out obvious meeting chatter.
_FILLER_PATTERNS = [
    r"\bmm-?hmm\b", r"\byou know\b", r"\bsort of\b", r"\bi mean\b", r"\bkind of\b",
    r"\bokay so\b", r"\bhmm+\b", r"\bum+\b", r"\buhh+\b", r"\buh+\b", r"\berm+\b",
    r"\byeah\b", r"\bbasically\b", r"\bliterally\b",
]


def _split_by_word_count(text: str, chunk_size: int = _MAX_SINGLE_PASS_WORDS) -> list[str]:
    words = text.split()
    return [" ".join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size)]


def _regex_clean(text: str) -> str:
    cleaned = text
    # Remove each filler along with an adjacent comma so we don't strand punctuation
    for pat in _FILLER_PATTERNS:
        cleaned = re.sub(r",?\s*" + pat + r"\s*,?", " ", cleaned, flags=re.IGNORECASE)
    # Collapse immediate duplicate words, allowing a comma between ("we, we")
    cleaned = re.sub(r"\b(\w+)(\s*,?\s*\1\b)+", r"\1", cleaned, flags=re.IGNORECASE)
    # Tidy debris: repeated commas, space-before-punctuation, leading commas
    cleaned = re.sub(r"(,\s*){2,}", ", ", cleaned)
    cleaned = re.sub(r"\s+([,.?!;:])", r"\1", cleaned)
    cleaned = re.sub(r"(?m)^\s*,\s*", "", cleaned)
    cleaned = re.sub(r"[ \t]{2,}", " ", cleaned)
    cleaned = "\n".join(line.strip(" ,") for line in cleaned.splitlines())
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def _llm_clean(text: str, visual_filenames: list[str]) -> str:
    visual_hint = ""
    if visual_filenames:
        visual_hint = f"\nVisual sources available in this meeting: {', '.join(visual_filenames)}"
    system = SystemMessage(content=CLEANING_PROMPT + visual_hint)
    llm = get_ingestion_llm()

    chunks = [text] if len(text.split()) <= _MAX_SINGLE_PASS_WORDS else _split_by_word_count(text)
    cleaned_parts = []
    for chunk in chunks:
        response = invoke_with_retry(llm, [system, HumanMessage(content=chunk)])
        cleaned_parts.append(response.content.strip())
    return "\n".join(cleaned_parts)


def _clean_text(text: str, visual_filenames: list[str]) -> str:
    if effective_config().cleaning_mode == "local":
        return _regex_clean(text)
    # llm mode: clean with the light model, falling back to regex if it is exhausted
    try:
        return _llm_clean(text, visual_filenames)
    except Exception:
        return _regex_clean(text)


def _process_transcript(doc: dict, visual_filenames: list[str]) -> str:
    text = get_storage().read(doc["storage_path"])
    return _clean_text(text, visual_filenames)


def _process_docx(doc: dict, visual_filenames: list[str]) -> str:
    import docx as python_docx
    storage = get_storage()
    raw_bytes = storage.read_bytes(doc["storage_path"])
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        tmp.write(raw_bytes)
        tmp_path = tmp.name
    try:
        document = python_docx.Document(tmp_path)
        text = "\n".join(p.text for p in document.paragraphs if p.text.strip())
    finally:
        os.remove(tmp_path)

    return _clean_text(text, visual_filenames)


def _process_pptx(doc: dict) -> str:
    from pptx import Presentation
    storage = get_storage()
    raw_bytes = storage.read_bytes(doc["storage_path"])
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp.write(raw_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_texts = []
        for i, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            parts.append(line)
            if parts:
                slide_texts.append(f"[Slide {i}]\n" + "\n".join(parts))
    finally:
        os.remove(tmp_path)
    return "\n\n".join(slide_texts)


def _process_pdf(doc: dict) -> str:
    import pymupdf4llm
    storage = get_storage()
    raw_bytes = storage.read_bytes(doc["storage_path"])
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        tmp.write(raw_bytes)
        tmp_path = tmp.name
    try:
        md = pymupdf4llm.to_markdown(tmp_path)
    finally:
        os.remove(tmp_path)
    return md


def _process_image(doc: dict) -> str:
    storage = get_storage()
    image_bytes = storage.read_bytes(doc["storage_path"])
    b64 = base64.b64encode(image_bytes).decode("utf-8")

    ext = doc["filename"].rsplit(".", 1)[-1].lower()
    mime_map = {"jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png",
                "webp": "image/webp"}
    mime = mime_map.get(ext, "image/png")

    llm = get_vision_llm()
    message = HumanMessage(content=[
        {"type": "text", "text": VISION_PROMPT},
        {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
    ])
    response = invoke_with_retry(llm, [message])
    return response.content.strip()


def _build_processed_text(doc: dict, visual_filenames: list[str]) -> str:
    ft = doc["file_type"]
    if ft == "transcript":
        return _process_transcript(doc, visual_filenames)
    if ft == "docx":
        return _process_docx(doc, visual_filenames)
    if ft == "pdf":
        return _process_pdf(doc)
    if ft == "pptx":
        return _process_pptx(doc)
    if ft == "image":
        return _process_image(doc)
    return ""


def _processed_text(doc: dict, visual_filenames: list[str]) -> str:
    """Return a document's processed text, reusing the cached cleaned_text if
    present, otherwise building it once and persisting it for future runs."""
    db = SessionLocal()
    try:
        record = db.query(Document).filter_by(id=doc["id"]).first()
        if record and record.cleaned_text:
            return record.cleaned_text
    finally:
        db.close()

    text = _build_processed_text(doc, visual_filenames)

    db = SessionLocal()
    try:
        record = db.query(Document).filter_by(id=doc["id"]).first()
        if record:
            record.cleaned_text = text
            db.commit()
    finally:
        db.close()
    return text


def ingestion_node(state: PipelineState) -> PipelineState:
    db = SessionLocal()
    try:
        run = db.query(PipelineRun).filter_by(id=state["run_id"]).first()
        statuses = dict(run.stage_statuses)
        statuses["ingestion"] = "running"
        run.stage_statuses = statuses
        run.current_stage = "ingestion"
        db.commit()
    finally:
        db.close()

    try:
        documents = state["documents"]
        visual_filenames = [d["filename"] for d in documents
                            if d["file_type"] in ("image", "pdf", "pptx")]

        transcript_parts = []
        visual_parts = []

        for doc in documents:
            text = _processed_text(doc, visual_filenames)
            if doc["file_type"] in ("transcript", "docx"):
                transcript_parts.append(text)
            else:
                visual_parts.append((doc["filename"], text))

        cleaned_content = "\n\n".join(transcript_parts)
        if visual_parts:
            cleaned_content += "\n\n--- VISUAL SOURCES ---\n\n"
            for filename, content in visual_parts:
                cleaned_content += f"[{filename}]\n{content}\n\n"

        db = SessionLocal()
        try:
            run = db.query(PipelineRun).filter_by(id=state["run_id"]).first()
            output = GeneratedOutput(
                project_id=state["project_id"],
                run_id=state["run_id"],
                stage="ingestion",
                content=cleaned_content,
            )
            db.add(output)
            statuses = dict(run.stage_statuses)
            statuses["ingestion"] = "complete"
            run.stage_statuses = statuses
            db.commit()
        finally:
            db.close()

        return {**state, "cleaned_content": cleaned_content, "current_stage": "ingestion"}

    except Exception as e:
        errors = list(state.get("errors", []))
        errors.append(f"Ingestion failed: {str(e)}")
        db = SessionLocal()
        try:
            run = db.query(PipelineRun).filter_by(id=state["run_id"]).first()
            project = db.query(Project).filter_by(id=state["project_id"]).first()
            statuses = dict(run.stage_statuses)
            statuses["ingestion"] = "failed"
            run.stage_statuses = statuses
            run.current_stage = "failed"
            run.error = str(e)
            project.status = "failed"
            db.commit()
        finally:
            db.close()
        raise
